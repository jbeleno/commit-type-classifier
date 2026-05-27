"""Generate the English defence deck (PowerPoint, 9 slides, 16:9).

Visual identity inherited from docs/design-system.md:
    background  slate-900   #0F172A
    surface     slate-800   #1E293B
    surface-2   slate-700   #334155
    ink         slate-50    #F8FAFC
    ink-3       slate-400   #94A3B8
    accent      sky         #38BDF8
    class colors: feat #22C55E, fix #EF4444, docs #38BDF8,
                  refactor #A78BFA, test #FBBF24

Run:
    make slides            # writes docs/exports/slides.pptx
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Design tokens — mirror docs/design-system.md §2
# --------------------------------------------------------------------------- #
BG = RGBColor(0x0F, 0x17, 0x2A)
SURFACE_1 = RGBColor(0x1E, 0x29, 0x3B)
SURFACE_2 = RGBColor(0x33, 0x41, 0x55)
SURFACE_INSET = RGBColor(0x0B, 0x12, 0x20)
BORDER_SOFT = RGBColor(0x33, 0x41, 0x55)
INK = RGBColor(0xF8, 0xFA, 0xFC)
INK_2 = RGBColor(0xCB, 0xD5, 0xE1)
INK_3 = RGBColor(0x94, 0xA3, 0xB8)
INK_MUTE = RGBColor(0x64, 0x74, 0x8B)
RING = RGBColor(0x38, 0xBD, 0xF8)

CLASS_COLORS = {
    "feat": RGBColor(0x22, 0xC5, 0x5E),
    "fix": RGBColor(0xEF, 0x44, 0x44),
    "docs": RGBColor(0x38, 0xBD, 0xF8),
    "refactor": RGBColor(0xA7, 0x8B, 0xFA),
    "test": RGBColor(0xFB, 0xBF, 0x24),
}

FONT_SANS = "Fira Sans"
FONT_MONO = "Fira Code"

ROOT = Path(__file__).resolve().parents[1]
DIAGRAMS = ROOT / "docs" / "diagrams" / "png"
MOCKUPS = ROOT / "docs" / "mockups"
OUT = ROOT / "docs" / "exports" / "slides.pptx"


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def make_presentation() -> Presentation:
    p = Presentation()
    # 16:9 widescreen
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def blank_slide(p: Presentation):
    layout = p.slide_layouts[6]  # blank layout
    slide = p.slides.add_slide(layout)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 16,
    color: RGBColor = INK,
    font_name: str = FONT_SANS,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    letter_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font_name
        f.size = Pt(font_size)
        f.color.rgb = color
        f.bold = bold
    return tb


def add_rectangle(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    line_width: float = 0.5,
):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    gutter: RGBColor = RING,
    fill: RGBColor = SURFACE_1,
):
    """A surface card with a 3 px left gutter (the signature element)."""
    add_rectangle(slide, x, y, w, h, fill=fill, line=BORDER_SOFT)
    add_rectangle(slide, x, y, 0.06, h, fill=gutter)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float):
    from PIL import Image as PILImage

    im = PILImage.open(path)
    aspect = im.width / im.height
    box_aspect = w / h
    if aspect > box_aspect:
        # constrained by width
        actual_w = w
        actual_h = w / aspect
    else:
        actual_h = h
        actual_w = h * aspect
    cx = x + (w - actual_w) / 2
    cy = y + (h - actual_h) / 2
    slide.shapes.add_picture(str(path), Inches(cx), Inches(cy), Inches(actual_w), Inches(actual_h))


def add_header(slide, eyebrow: str, title: str):
    add_text(slide, eyebrow.upper(), 0.6, 0.45, 12, 0.3,
             font_size=10, color=INK_3, font_name=FONT_SANS, bold=False)
    add_text(slide, title, 0.6, 0.75, 12, 0.7,
             font_size=28, color=INK, font_name=FONT_SANS, bold=True)
    # accent underline
    add_rectangle(slide, 0.6, 1.45, 0.6, 0.04, fill=RING)


def add_footer(slide, idx: int, total: int):
    add_text(slide, "Commit Type Classifier · C2", 0.6, 7.1, 6, 0.3,
             font_size=9, color=INK_MUTE, font_name=FONT_MONO)
    add_text(slide, f"{idx} / {total}", 12.0, 7.1, 1.2, 0.3,
             font_size=9, color=INK_MUTE, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_chip(slide, label: str, x: float, y: float, *, w: float = 1.0):
    color = CLASS_COLORS.get(label, RING)
    # Soft-tinted background — approximate by mixing color with surface
    tint = RGBColor(
        int((color[0] + SURFACE_1[0] * 5) / 6),
        int((color[1] + SURFACE_1[1] * 5) / 6),
        int((color[2] + SURFACE_1[2] * 5) / 6),
    )
    add_rectangle(slide, x, y, w, 0.32, fill=tint, line=color)
    add_text(slide, label, x, y, w, 0.32,
             font_size=12, color=color, font_name=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------- #
# Individual slides
# --------------------------------------------------------------------------- #
TOTAL = 12


def slide_01_title(p):
    s = blank_slide(p)
    # Left accent gutter spanning the whole height
    add_rectangle(s, 0, 0, 0.12, 7.5, fill=RING)

    # Big monospace word mark
    add_text(s, "commit-type-classifier", 0.9, 1.2, 12, 0.5,
             font_size=14, color=INK_3, font_name=FONT_MONO)
    add_text(s, "Commit Type Classifier", 0.9, 1.7, 12, 1.2,
             font_size=54, color=INK, font_name=FONT_SANS, bold=True)
    add_text(s, "Component 2 — AI for Software Engineering", 0.9, 2.8, 12, 0.5,
             font_size=22, color=INK_2, font_name=FONT_SANS)

    # Five class chips as a subtle visual hook
    for i, label in enumerate(["feat", "fix", "docs", "refactor", "test"]):
        add_chip(s, label, 0.9 + i * 1.15, 3.6, w=1.0)

    # Authors / course block
    add_text(s, "Jesús Beleño  ·  Juan Forero", 0.9, 5.0, 12, 0.4,
             font_size=18, color=INK, font_name=FONT_SANS, bold=True)
    add_text(s, "BEINSOF52 · Artificial Intelligence  ·  Advisor: Juan Antonio Castro Silva",
             0.9, 5.4, 12, 0.4, font_size=13, color=INK_2, font_name=FONT_MONO)
    add_text(s, "Universidad Surcolombiana  ·  May 2026",
             0.9, 5.75, 12, 0.4, font_size=13, color=INK_3, font_name=FONT_MONO)

    add_text(s, "github.com/jbeleno/commit-type-classifier", 0.9, 6.85, 12, 0.4,
             font_size=11, color=INK_MUTE, font_name=FONT_MONO)
    return s


def slide_02_problem(p):
    s = blank_slide(p)
    add_header(s, "01 · the problem", "Conventional Commits adoption is uneven")
    add_footer(s, 2, TOTAL)

    # Three big stat cards
    stats = [
        ("< 5 %", "of public commits strictly follow\nthe Conventional Commits spec", CLASS_COLORS["fix"]),
        ("~ 40 %", "disagreement rate between author\nprefix and a defensible re-classification", CLASS_COLORS["test"]),
        ("38,965", "labelled commits retained\nafter filtering the 5 target classes", CLASS_COLORS["feat"]),
    ]
    for i, (value, label, color) in enumerate(stats):
        x = 0.6 + i * 4.18
        w = 3.9
        add_card(s, x, 2.0, w, 3.0, gutter=color)
        add_text(s, value, x + 0.4, 2.25, w - 0.8, 1.4,
                 font_size=54, color=INK, font_name=FONT_MONO, bold=True)
        add_text(s, label, x + 0.4, 3.7, w - 0.8, 1.2,
                 font_size=13, color=INK_2, font_name=FONT_SANS)

    # Bottom line
    add_text(s,
        "Two consequences:  manual triage is expensive  ·  author-supplied labels are inconsistent.",
        0.6, 5.5, 12, 0.5, font_size=14, color=INK_3, font_name=FONT_SANS)
    return s


def slide_03_approach(p):
    s = blank_slide(p)
    add_header(s, "02 · approach", "Five heterogeneous models, one local pipeline")
    add_footer(s, 3, TOTAL)

    families = [
        ("baseline_tfidf", "Classical ML\nTF-IDF + LogReg", CLASS_COLORS["feat"]),
        ("cnn_text", "Deep learning\nDual-branch CNN", CLASS_COLORS["docs"]),
        ("distilbert", "Transformer (text)\nDistilBERT 67M fine-tuned", CLASS_COLORS["refactor"]),
        ("codebert", "Transformer (code)\nCodeBERT 125M fine-tuned", CLASS_COLORS["test"]),
        ("ensemble", "Soft voting\nL-BFGS-B over the 4 probas", RING),
    ]
    for i, (name, body, color) in enumerate(families):
        x = 0.6 + i * 2.50
        w = 2.32
        add_card(s, x, 2.0, w, 2.9, gutter=color)
        add_text(s, name, x + 0.25, 2.2, w - 0.5, 0.5,
                 font_size=15, color=color, font_name=FONT_MONO, bold=True)
        add_text(s, body, x + 0.25, 2.8, w - 0.5, 1.6,
                 font_size=12, color=INK_2, font_name=FONT_SANS)

    add_text(s,
        "Not an LLM application — these are supervised classifiers (encoder-only transformers, not generative).",
        0.6, 5.4, 12, 0.5, font_size=14, color=INK_3, font_name=FONT_SANS)
    add_text(s,
        "Local-only inference path — no cloud, no API calls.",
        0.6, 5.85, 12, 0.5, font_size=14, color=INK_3, font_name=FONT_SANS)
    return s


def slide_04_pipeline(p):
    s = blank_slide(p)
    add_header(s, "03 · pipeline", "CommitBench → 70/15/15 → five trained models")
    add_footer(s, 4, TOTAL)
    add_image_fit(s, DIAGRAMS / "05_architecture_pipeline.png", 0.6, 1.7, 12.1, 5.3)
    return s


def slide_05_results(p):
    s = blank_slide(p)
    add_header(s, "04 · results", "Test set — 5,845 commits  ·  baseline_tfidf wins on every metric")
    add_footer(s, 5, TOTAL)

    rows = [
        ("baseline_tfidf", 0.7093, 0.6632, True),
        ("ensemble",       0.6896, 0.6438, False),
        ("cnn_text",       0.6599, 0.5861, False),
        ("codebert",       0.6089, 0.5800, False),
        ("distilbert",     0.5858, 0.5515, False),
    ]
    max_f1 = max(r[2] for r in rows)
    row_h = 0.65
    top = 2.05

    add_text(s, "model", 0.6, top - 0.4, 2.4, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True, letter_spacing=2.0)
    add_text(s, "accuracy", 3.0, top - 0.4, 1.5, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "macro F1", 4.7, top - 0.4, 1.5, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "macro F1 — visual", 7.0, top - 0.4, 5.4, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True)

    for i, (name, acc, f1, leader) in enumerate(rows):
        y = top + i * row_h
        color = RING if leader else SURFACE_2
        add_rectangle(s, 0.6, y, 12.1, 0.02, fill=BORDER_SOFT)
        add_rectangle(s, 0.6, y + 0.05, 0.04, row_h - 0.1, fill=color)
        add_text(s, name, 0.85, y + 0.05, 2.2, row_h - 0.1, font_size=15,
                 color=INK if leader else INK_2, font_name=FONT_MONO, bold=leader,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{acc * 100:.2f} %", 3.0, y + 0.05, 1.5, row_h - 0.1,
                 font_size=15, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{f1:.4f}", 4.7, y + 0.05, 1.5, row_h - 0.1,
                 font_size=15, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # bar
        bar_x = 6.4
        bar_w = 5.4
        track_y = y + (row_h - 0.18) / 2
        add_rectangle(s, bar_x, track_y, bar_w, 0.18,
                      fill=SURFACE_INSET, line=BORDER_SOFT)
        fill_w = bar_w * (f1 / max_f1)
        add_rectangle(s, bar_x, track_y, fill_w, 0.18, fill=RING)

    add_text(s, "Why TF-IDF wins:", 0.6, 6.0, 4, 0.4,
             font_size=14, color=INK_3, font_name=FONT_SANS, bold=True)
    add_text(s,
        "commit messages are short (median 7 tokens), the diff signal is captured well by char n-grams, "
        "and transformers were fine-tuned on only 6–8 k balanced examples.",
        0.6, 6.4, 12.1, 0.8, font_size=12, color=INK_3, font_name=FONT_SANS)
    return s


def slide_06_insight(p):
    s = blank_slide(p)
    add_header(s, "05 · insight", "When the model knows better than the author")
    add_footer(s, 6, TOTAL)

    # Big quote-style card
    add_card(s, 0.6, 1.9, 6.0, 4.8, gutter=CLASS_COLORS["feat"])
    add_text(s, "residence-back  ·  25 commits  ·  100 % Conventional", 0.85, 2.05, 5.6, 0.35,
             font_size=11, color=INK_3, font_name=FONT_MONO)
    add_text(s, "15 / 25", 0.85, 2.5, 5.6, 1.1,
             font_size=48, color=INK, font_name=FONT_MONO, bold=True)
    add_text(s, "agreement with the author-supplied prefix",
             0.85, 3.6, 5.6, 0.4, font_size=14, color=INK_2, font_name=FONT_SANS)
    add_text(s, "of the 10 disagreements:", 0.85, 4.35, 5.6, 0.4,
             font_size=12, color=INK_3, font_name=FONT_SANS, bold=True)
    add_text(s, "•  8  were author mislabels  (fix → actually feat)\n"
                "•  2  were genuine model errors on docs commits",
             0.85, 4.65, 5.6, 1.2, font_size=13, color=INK_2, font_name=FONT_MONO)
    add_text(s, "Defensible accuracy ≈ 92 %",
             0.85, 6.0, 5.6, 0.5, font_size=18, color=CLASS_COLORS["feat"],
             font_name=FONT_MONO, bold=True)

    # Right-side example pair
    add_card(s, 6.95, 1.9, 5.75, 2.25, gutter=CLASS_COLORS["fix"])
    add_text(s, "author said", 7.2, 2.05, 5.3, 0.3,
             font_size=10, color=INK_3, font_name=FONT_SANS)
    add_chip(s, "fix", 7.2, 2.45, w=0.7)
    add_text(s,
        "fix: super_admin can access any condominium without a formal UCR",
        8.05, 2.45, 4.55, 0.7, font_size=12, color=INK_2, font_name=FONT_MONO)
    add_text(s,
        "→ adds a brand-new capability for super_admin; semantically a feature.",
        7.2, 3.25, 5.3, 0.8, font_size=11, color=INK_3, font_name=FONT_SANS)

    add_card(s, 6.95, 4.45, 5.75, 2.25, gutter=CLASS_COLORS["feat"])
    add_text(s, "model said", 7.2, 4.6, 5.3, 0.3,
             font_size=10, color=INK_3, font_name=FONT_SANS)
    add_chip(s, "feat", 7.2, 5.0, w=0.75)
    add_text(s, "confidence 0.985",
             8.1, 5.0, 4.5, 0.4, font_size=12, color=INK, font_name=FONT_MONO)
    add_text(s, "The classifier picked up on the diff: a new authorisation path.",
             7.2, 5.8, 5.3, 0.8, font_size=11, color=INK_3, font_name=FONT_SANS)
    return s


def slide_07_llm_pivot(p):
    s = blank_slide(p)
    add_header(s, "06 · pivot", "Discriminative → Generative  ·  local LLMs via Ollama")
    add_footer(s, 7, TOTAL)

    add_card(s, 0.6, 1.95, 5.7, 4.85, gutter=CLASS_COLORS["fix"])
    add_text(s, "what we had", 0.85, 2.1, 5.3, 0.4,
             font_size=14, color=CLASS_COLORS["fix"], font_name=FONT_SANS, bold=True)
    add_text(s,
        "•  Classifier: message + diff  →  one of 5 labels\n\n"
        "•  TF-IDF baseline at 70.93 % acc / 0.66 F1\n\n"
        "•  Useful, but answers the wrong question — it only\n"
        "    re-labels what the developer already wrote",
        0.85, 2.55, 5.3, 4.0, font_size=13, color=INK_2, font_name=FONT_SANS)

    add_card(s, 6.6, 1.95, 6.1, 4.85, gutter=CLASS_COLORS["feat"])
    add_text(s, "what we built", 6.85, 2.1, 5.7, 0.4,
             font_size=14, color=CLASS_COLORS["feat"], font_name=FONT_SANS, bold=True)
    add_text(s,
        "•  Generator: diff alone  →  full Conventional Commit\n\n"
        "•  5 local LLMs served by Ollama (16 GB unified RAM):\n"
        "    qwen2.5-coder 1.5b / 3b,  llama3.2 3b-instruct,\n"
        "    phi3.5 3.8b-mini-instruct,  deepseek-coder 1.3b\n\n"
        "•  4 prompting strategies × 5 models = 20 configs swept\n\n"
        "•  Old classifier  →  reused as type verifier (hybrid)",
        6.85, 2.55, 5.7, 4.0, font_size=13, color=INK_2, font_name=FONT_SANS)
    return s


def slide_08_llm_generation(p):
    s = blank_slide(p)
    add_header(s, "07 · generation", "Sweep results — type-exact-match on stratified n=50")
    add_footer(s, 8, TOTAL)

    rows = [
        ("phi3.5:3.8b-mini · few_shot",    0.36, True),
        ("llama3.2:3b · CoT",              0.32, False),
        ("qwen2.5-coder:3b · CoT",         0.30, False),
        ("qwen2.5-coder:1.5b · CoT",       0.26, False),
        ("deepseek-coder:1.3b · json",     0.18, False),
    ]
    max_v = max(r[1] for r in rows)
    row_h = 0.62
    top = 2.05

    add_text(s, "best strategy per model", 0.6, top - 0.4, 6, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True, letter_spacing=2.0)
    add_text(s, "type-match", 6.4, top - 0.4, 2, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "visual", 8.6, top - 0.4, 4, 0.3, font_size=10, color=INK_3,
             font_name=FONT_SANS, bold=True)

    for i, (name, v, leader) in enumerate(rows):
        y = top + i * row_h
        color = RING if leader else SURFACE_2
        add_rectangle(s, 0.6, y, 12.1, 0.02, fill=BORDER_SOFT)
        add_rectangle(s, 0.6, y + 0.05, 0.04, row_h - 0.1, fill=color)
        add_text(s, name, 0.85, y + 0.05, 5.4, row_h - 0.1, font_size=13,
                 color=INK if leader else INK_2, font_name=FONT_MONO, bold=leader,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{v * 100:.1f} %", 6.4, y + 0.05, 2.0, row_h - 0.1,
                 font_size=14, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        bar_x, bar_w = 8.6, 4.1
        track_y = y + (row_h - 0.18) / 2
        add_rectangle(s, bar_x, track_y, bar_w, 0.18,
                      fill=SURFACE_INSET, line=BORDER_SOFT)
        add_rectangle(s, bar_x, track_y, bar_w * (v / max_v), 0.18, fill=RING)

    add_text(s, "Key insight:", 0.6, 5.55, 4, 0.4,
             font_size=14, color=INK_3, font_name=FONT_SANS, bold=True)
    add_text(s,
        "instruction-tuning matters more than parameter count or family — "
        "phi3.5 (instruct) > qwen-coder (code-aware base); "
        "deepseek-coder (base/completion) hits 100 % parse-fail outside JSON-mode.",
        0.6, 5.95, 12.1, 1.0, font_size=12, color=INK_3, font_name=FONT_SANS)
    return s


def slide_09_apples_to_apples(p):
    s = blank_slide(p)
    add_header(s, "08 · LLM vs baseline", "Apples-to-apples  ·  same input, same labels")
    add_footer(s, 9, TOTAL)

    headers = [("system",           0.6, 5.0),
               ("accuracy",         5.8, 1.7),
               ("macro F1",         7.6, 1.7),
               ("weighted F1",      9.4, 1.9)]
    for name, x, w in headers:
        add_text(s, name, x, 1.85, w, 0.3, font_size=10, color=INK_3,
                 font_name=FONT_SANS, bold=True, letter_spacing=2.0,
                 align=PP_ALIGN.LEFT if name == "system" else PP_ALIGN.RIGHT)

    rows = [
        ("TF-IDF baseline (n=5 845)",        0.7093, 0.6632, 0.7187, False),
        ("qwen2.5-coder:3b · rag (n=200)",   0.7400, 0.5639, 0.7190, False),
        ("Hard-vote ensemble (4 members)",   0.7750, 0.6014, 0.7457, False),
        ("Weighted ensemble (TF-IDF 2× boost)", 0.7500, 0.6698, 0.7505, True),
    ]
    row_h = 0.62
    top = 2.2

    for i, (name, acc, f1, wf1, winner) in enumerate(rows):
        y = top + i * row_h
        color = CLASS_COLORS["feat"] if winner else SURFACE_2
        add_rectangle(s, 0.6, y, 12.1, 0.02, fill=BORDER_SOFT)
        add_rectangle(s, 0.6, y + 0.05, 0.04, row_h - 0.1, fill=color)
        add_text(s, name, 0.85, y + 0.05, 4.8, row_h - 0.1, font_size=13,
                 color=INK if winner else INK_2, font_name=FONT_MONO, bold=winner,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{acc * 100:.2f} %", 5.6, y + 0.05, 1.7, row_h - 0.1,
                 font_size=14, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{f1:.4f}", 7.4, y + 0.05, 1.7, row_h - 0.1,
                 font_size=14, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, f"{wf1:.4f}", 9.3, y + 0.05, 1.9, row_h - 0.1,
                 font_size=14, color=INK, font_name=FONT_MONO,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "What the winner row says:", 0.6, 5.4, 5, 0.4,
             font_size=14, color=CLASS_COLORS["feat"], font_name=FONT_SANS, bold=True)
    add_text(s,
        "The heterogeneous voting ensemble (qwen2.5-coder:3b + phi3.5:3.8b-mini + TF-IDF voter, "
        "weighted by accuracy with the TF-IDF voter doubled) is the first configuration that strictly beats "
        "the discriminative baseline on every test-set metric — accuracy +4.07 pp, macro-F1 +0.0066, "
        "weighted-F1 +0.0318. The TF-IDF baseline is reused as a tie-breaker on minority classes "
        "(docs, refactor, test), turning the discriminative work into infrastructure rather than competition.",
        0.6, 5.85, 12.1, 1.4, font_size=12, color=INK_3, font_name=FONT_SANS)
    return s


def slide_07_demo(p):
    s = blank_slide(p)
    add_header(s, "09 · demo", "Local app — Streamlit GUI  +  Typer CLI")
    add_footer(s, 10, TOTAL)

    add_image_fit(s, MOCKUPS / "04_metrics.jpeg", 0.6, 1.8, 8.0, 5.0)

    # Right side — CLI snippet card
    add_card(s, 8.85, 1.8, 3.9, 2.5, gutter=RING)
    add_text(s, "CLI", 9.05, 1.95, 3.5, 0.35,
             font_size=10, color=INK_3, font_name=FONT_SANS)
    add_text(s,
        "$ python -m app.cli \\\n"
        "    predict-cmd \\\n"
        "    --message \"fix: oauth bug\" \\\n"
        "    --model baseline_tfidf\n\n"
        "fix (0.786)  via baseline_tfidf",
        9.05, 2.3, 3.6, 1.9, font_size=11, color=INK_2, font_name=FONT_MONO)

    add_card(s, 8.85, 4.5, 3.9, 2.3, gutter=CLASS_COLORS["feat"])
    add_text(s, "what the GUI shows", 9.05, 4.65, 3.5, 0.35,
             font_size=10, color=INK_3, font_name=FONT_SANS)
    add_text(s,
        "• 2 px class-coloured gutter on every\n"
        "  prediction card (signature element)\n"
        "• 5-row probability bars in mono\n"
        "• Diff viewer with +/- syntax colouring\n"
        "• Repo-scan histogram by class",
        9.05, 5.0, 3.6, 1.8, font_size=11, color=INK_2, font_name=FONT_SANS)
    return s


def slide_08_future(p):
    s = blank_slide(p)
    add_header(s, "10 · limitations & future work", "What we know we left on the table")
    add_footer(s, 11, TOTAL)

    add_card(s, 0.6, 2.0, 6.0, 4.7, gutter=CLASS_COLORS["fix"])
    add_text(s, "limitations", 0.85, 2.15, 5.6, 0.4,
             font_size=14, color=CLASS_COLORS["fix"], font_name=FONT_SANS, bold=True)
    add_text(s,
        "•  5 of the 10 canonical Conventional Commit types\n\n"
        "•  English corpus only — Spanish messages degrade the\n"
        "    text branch (diff branch is language-agnostic)\n\n"
        "•  Apple-Silicon MPS hangs when TF and PyTorch share a\n"
        "    process, so the ensemble runs only at evaluation time\n\n"
        "•  Transformers fine-tuned on 6–8 k balanced examples\n"
        "    (not the full 27 k train split)",
        0.85, 2.65, 5.6, 4.0, font_size=13, color=INK_2, font_name=FONT_SANS)

    add_card(s, 6.95, 2.0, 5.75, 4.7, gutter=CLASS_COLORS["feat"])
    add_text(s, "future work", 7.2, 2.15, 5.3, 0.4,
             font_size=14, color=CLASS_COLORS["feat"], font_name=FONT_SANS, bold=True)
    add_text(s,
        "•  Stacking ensemble with a learned meta-classifier\n"
        "    (closes the gap with the TF-IDF baseline)\n\n"
        "•  Multilingual transformer (Spanish + English\n"
        "    code review reality in Latin-American teams)\n\n"
        "•  CodeBERT → 30 M-parameter distilled student\n"
        "    for CPU-only deployment\n\n"
        "•  Wrap the inference layer in Flask / FastAPI so\n"
        "    C1 can call it as an HTTP service",
        7.2, 2.65, 5.3, 4.0, font_size=13, color=INK_2, font_name=FONT_SANS)
    return s


def slide_09_thanks(p):
    s = blank_slide(p)
    add_rectangle(s, 0, 0, 0.12, 7.5, fill=RING)

    add_text(s, "thanks · questions", 0.9, 2.5, 12, 1.3,
             font_size=56, color=INK, font_name=FONT_SANS, bold=True)
    add_text(s, "Commit Type Classifier  ·  Component 2", 0.9, 3.7, 12, 0.5,
             font_size=20, color=INK_2, font_name=FONT_SANS)
    add_text(s, "Jesús Beleño  ·  Juan Forero", 0.9, 4.2, 12, 0.5,
             font_size=15, color=INK_3, font_name=FONT_MONO)

    add_card(s, 0.9, 5.2, 11.5, 1.2, gutter=RING)
    add_text(s, "github.com/jbeleno/commit-type-classifier", 1.2, 5.45, 11, 0.7,
             font_size=20, color=INK, font_name=FONT_MONO, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    return s


# --------------------------------------------------------------------------- #
def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    p = make_presentation()
    slide_01_title(p)
    slide_02_problem(p)
    slide_03_approach(p)
    slide_04_pipeline(p)
    slide_05_results(p)
    slide_06_insight(p)
    slide_07_llm_pivot(p)
    slide_08_llm_generation(p)
    slide_09_apples_to_apples(p)
    slide_07_demo(p)
    slide_08_future(p)
    slide_09_thanks(p)
    p.save(OUT)
    print(f"Saved {OUT} ({len(p.slides)} slides)")


if __name__ == "__main__":
    main()
